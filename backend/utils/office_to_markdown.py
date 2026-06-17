"""
Office 文档原生转 Markdown 转换器

不依赖 MarkItDown，使用 python-docx / python-pptx / openpyxl 直接解析文档，
保留图片在原文中的位置并上传到 RustFS。

支持格式:
  - DOCX  → python-docx
  - PPTX  → python-pptx
  - XLSX  → openpyxl
"""

import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from loguru import logger


def _strip_external_links_to_tmp(xlsx_path: str) -> str:
    """
    复制一份 xlsx，移除 xl/externalLinks/ 及 workbook.xml 中的 <externalReferences> 节点。

    openpyxl 对含外部工作簿引用（externalReferences，常见于 WPS / 链接了其它表格的 Excel）
    的 xlsx 解析有兼容缺陷：ExternalReference.__init__() missing 'id'。移除这些外部链接后
    即可正常加载，对表格数据本身无影响。返回处理后的临时文件路径（调用方负责删除）。
    """
    import os
    import zipfile
    import tempfile

    fd, tmp_path = tempfile.mkstemp(suffix=".xlsx")
    os.close(fd)
    with zipfile.ZipFile(xlsx_path, "r") as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            if item.startswith("xl/externalLinks/"):
                continue  # 丢弃外部链接定义
            data = zin.read(item)
            if item == "xl/workbook.xml":
                text = data.decode("utf-8", errors="ignore")
                text = re.sub(r"<externalReferences>.*?</externalReferences>", "", text, flags=re.DOTALL)
                text = re.sub(r"<externalReferences\s*/>", "", text)
                data = text.encode("utf-8")
            elif item == "xl/_rels/workbook.xml.rels":
                text = data.decode("utf-8", errors="ignore")
                text = re.sub(r"<Relationship[^>]*externalLink[^>]*?/>", "", text)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    return tmp_path


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def docx_to_markdown(
    docx_path: str,
    images_dir: str,
    rustfs_client=None,
) -> Tuple[str, List[Dict]]:
    """
    将 DOCX 转换为 Markdown，图片保留原始位置。

    Returns:
        (markdown_content, images)
        images: [{"src": filename, "alt": alt, "local_path": ..., "rustfs_url": ...}]
    """
    from docx import Document
    from docx.oxml.ns import qn

    images_path = Path(images_dir)
    images_path.mkdir(parents=True, exist_ok=True)

    doc = Document(docx_path)
    images: List[Dict] = []
    image_counter = 0
    output_parts: List[str] = []

    def _upload(img_bytes: bytes, img_ext: str, img_name: str) -> Tuple[str, Optional[str]]:
        local_path = images_path / img_name
        local_path.write_bytes(img_bytes)
        rustfs_url = None
        if rustfs_client:
            try:
                rustfs_url = rustfs_client.upload_file(str(local_path))
            except Exception as e:
                logger.warning(f"⚠️  RustFS upload failed for {img_name}: {e}")
        return str(local_path), rustfs_url

    def _extract_inline_images(element) -> List[str]:
        """从 XML element 中提取所有内联图片，返回对应的 <img> 或 markdown 字符串列表"""
        nonlocal image_counter
        result = []
        for drawing in element.iter(qn("w:drawing")):
            blip = drawing.find(".//" + qn("a:blip"))
            if blip is None:
                continue
            r_embed = blip.get(qn("r:embed"))
            if not r_embed:
                continue
            try:
                img_part = doc.part.related_parts[r_embed]
            except KeyError:
                continue
            img_bytes = img_part.blob
            content_type = img_part.content_type  # e.g. image/png
            img_ext = content_type.split("/")[-1].lower()
            if img_ext == "jpeg":
                img_ext = "jpg"
            # 过滤非图片（例如 wmf/emf 矢量格式，Pillow 通常无法处理）
            if img_ext not in ("png", "jpg", "jpeg", "gif", "bmp", "webp", "tiff"):
                logger.debug(f"Skipping unsupported image type: {img_ext}")
                continue
            image_counter += 1
            img_name = f"image{image_counter}.{img_ext}"
            local_path, rustfs_url = _upload(img_bytes, img_ext, img_name)
            img_info = {
                "src": img_name,
                "alt": img_name,
                "local_path": local_path,
                "rustfs_url": rustfs_url,
            }
            images.append(img_info)
            if rustfs_url:
                result.append(f'<img src="{rustfs_url}" alt="{img_name}">')
            else:
                result.append(f"![{img_name}](images/{img_name})")
        return result

    def _heading_prefix(para) -> str:
        style = para.style.name if para.style else ""
        mapping = {
            "Heading 1": "# ",
            "标题 1": "# ",
            "Heading 2": "## ",
            "标题 2": "## ",
            "Heading 3": "### ",
            "标题 3": "### ",
            "Heading 4": "#### ",
            "标题 4": "#### ",
            "Heading 5": "##### ",
            "标题 5": "##### ",
            "Heading 6": "###### ",
            "标题 6": "###### ",
        }
        for key, prefix in mapping.items():
            if style.startswith(key):
                return prefix
        return ""

    def _table_to_markdown(table) -> str:
        rows = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append(cells)
        if not rows:
            return ""
        col_count = max(len(r) for r in rows)
        # 补齐列数
        rows = [r + [""] * (col_count - len(r)) for r in rows]
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * col_count) + " |"
        body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
        return "\n".join(filter(None, [header, separator, body]))

    # 按 body 顺序遍历段落和表格

    body = doc.element.body
    para_idx = 0
    table_idx = 0

    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            # 段落
            if para_idx < len(doc.paragraphs):
                para = doc.paragraphs[para_idx]
            else:
                para_idx += 1
                continue
            para_idx += 1

            prefix = _heading_prefix(para)
            text = para.text.strip()
            img_tags = _extract_inline_images(child)

            if img_tags:
                if text:
                    output_parts.append(f"{prefix}{text}")
                output_parts.extend(img_tags)
            elif text:
                output_parts.append(f"{prefix}{text}")

        elif tag == "tbl":
            # 表格
            if table_idx < len(doc.tables):
                table = doc.tables[table_idx]
                md_table = _table_to_markdown(table)
                if md_table:
                    output_parts.append(md_table)
            table_idx += 1

    return "\n\n".join(output_parts), images


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def pptx_to_markdown(
    pptx_path: str,
    images_dir: str,
    rustfs_client=None,
) -> Tuple[str, List[Dict]]:
    """
    将 PPTX 转换为 Markdown，图片保留在对应幻灯片内。

    Returns:
        (markdown_content, images)
    """
    from pptx import Presentation

    PICTURE_TYPE = 13  # MSO_SHAPE_TYPE.PICTURE

    images_path = Path(images_dir)
    images_path.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    images: List[Dict] = []
    image_counter = 0
    output_parts: List[str] = []

    def _upload(img_bytes: bytes, img_ext: str, img_name: str) -> Tuple[str, Optional[str]]:
        local_path = images_path / img_name
        local_path.write_bytes(img_bytes)
        rustfs_url = None
        if rustfs_client:
            try:
                rustfs_url = rustfs_client.upload_file(str(local_path))
            except Exception as e:
                logger.warning(f"⚠️  RustFS upload failed for {img_name}: {e}")
        return str(local_path), rustfs_url

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: List[str] = [f"## 幻灯片 {slide_num}"]

        # 按位置排序（从上到下、从左到右）
        shapes = sorted(
            slide.shapes,
            key=lambda s: (s.top if s.top is not None else 0, s.left if s.left is not None else 0),
        )

        for shape in shapes:
            # 图片
            if shape.shape_type == PICTURE_TYPE:
                image_counter += 1
                img = shape.image
                img_ext = img.ext.lower()
                if img_ext == "jpeg":
                    img_ext = "jpg"
                img_name = f"slide{slide_num}_image{image_counter}.{img_ext}"
                local_path, rustfs_url = _upload(img.blob, img_ext, img_name)
                img_info = {
                    "src": img_name,
                    "alt": shape.name or img_name,
                    "local_path": local_path,
                    "rustfs_url": rustfs_url,
                }
                images.append(img_info)
                if rustfs_url:
                    slide_parts.append(f'<img src="{rustfs_url}" alt="{img_info["alt"]}">')
                else:
                    slide_parts.append(f'![{img_info["alt"]}](images/{img_name})')
                continue

            # 文本框 / 占位符
            if shape.has_text_frame:
                text_lines = []
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        text_lines.append(line)
                if text_lines:
                    slide_parts.append("\n".join(text_lines))

            # 表格
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    col_count = max(len(r) for r in rows)
                    rows = [r + [""] * (col_count - len(r)) for r in rows]
                    md_table = "| " + " | ".join(rows[0]) + " |\n"
                    md_table += "| " + " | ".join(["---"] * col_count) + " |\n"
                    md_table += "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
                    slide_parts.append(md_table)

        output_parts.append("\n\n".join(slide_parts))

    return "\n\n---\n\n".join(output_parts), images


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def xlsx_to_markdown(
    xlsx_path: str,
    images_dir: str,
    rustfs_client=None,
) -> Tuple[str, List[Dict]]:
    """
    将 XLSX 转换为 Markdown。
    表格内容转为 Markdown 表格，嵌入图片追加在所属 Sheet 末尾。
    （Excel 图片是浮动对象，无法精确定位到单元格，追加是最可靠的方案。）

    Returns:
        (markdown_content, images)
    """
    from openpyxl import load_workbook

    images_path = Path(images_dir)
    images_path.mkdir(parents=True, exist_ok=True)

    try:
        wb = load_workbook(xlsx_path, data_only=True)
    except TypeError as e:
        # openpyxl 对含 externalReferences（外部工作簿引用）的 xlsx 解析有兼容缺陷
        # （ExternalReference.__init__() missing 'id'）。预处理移除外部链接后重试。
        if "ExternalReference" in str(e) or "external" in str(e).lower():
            logger.warning(f"⚠️  xlsx 含外部引用导致 openpyxl 解析失败，移除外部链接后重试: {e}")
            cleaned = _strip_external_links_to_tmp(xlsx_path)
            try:
                wb = load_workbook(cleaned, data_only=True)
            finally:
                try:
                    Path(cleaned).unlink()
                except Exception:
                    pass
        else:
            raise
    images: List[Dict] = []
    image_counter = 0
    output_parts: List[str] = []

    def _upload(img_bytes: bytes, img_ext: str, img_name: str) -> Tuple[str, Optional[str]]:
        local_path = images_path / img_name
        local_path.write_bytes(img_bytes)
        rustfs_url = None
        if rustfs_client:
            try:
                rustfs_url = rustfs_client.upload_file(str(local_path))
            except Exception as e:
                logger.warning(f"⚠️  RustFS upload failed for {img_name}: {e}")
        return str(local_path), rustfs_url

    for ws in wb.worksheets:
        sheet_parts: List[str] = [f"## {ws.title}"]

        # 表格数据
        max_row = ws.max_row or 0
        max_col = ws.max_column or 0
        if max_row > 0 and max_col > 0:
            rows = list(
                ws.iter_rows(
                    min_row=1,
                    max_row=min(max_row, 2000),
                    values_only=True,
                )
            )
            # 过滤全空行
            rows = [r for r in rows if any(c is not None for c in r)]
            if rows:

                def _cell(v):
                    if v is None:
                        return ""
                    s = str(v).replace("|", "\\|").replace("\n", " ")
                    return s

                header = [_cell(c) for c in rows[0]]
                sheet_parts.append("| " + " | ".join(header) + " |")
                sheet_parts.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in rows[1:]:
                    cells = [_cell(c) for c in row]
                    # 补齐列数
                    while len(cells) < len(header):
                        cells.append("")
                    sheet_parts.append("| " + " | ".join(cells) + " |")

        # 嵌入图片（浮动，追加到 sheet 末尾）
        sheet_image_tags: List[str] = []
        if hasattr(ws, "_images") and ws._images:
            for img_obj in ws._images:
                image_counter += 1
                try:
                    img_bytes = img_obj._data()
                    # openpyxl 的 Image 对象有 format 属性
                    img_fmt = getattr(img_obj, "format", None) or "png"
                    img_ext = img_fmt.lower()
                    if img_ext == "jpeg":
                        img_ext = "jpg"
                    img_name = f"{ws.title}_image{image_counter}.{img_ext}"
                    local_path, rustfs_url = _upload(img_bytes, img_ext, img_name)
                    img_info = {
                        "src": img_name,
                        "alt": img_name,
                        "local_path": local_path,
                        "rustfs_url": rustfs_url,
                    }
                    images.append(img_info)
                    if rustfs_url:
                        sheet_image_tags.append(f'<img src="{rustfs_url}" alt="{img_name}">')
                    else:
                        sheet_image_tags.append(f"![{img_name}](images/{img_name})")
                except Exception as e:
                    logger.warning(f"⚠️  Failed to extract image from {ws.title}: {e}")

        if sheet_image_tags:
            sheet_parts.append("\n".join(sheet_image_tags))

        output_parts.append("\n\n".join(sheet_parts))

    return "\n\n---\n\n".join(output_parts), images


# ---------------------------------------------------------------------------
# 统一入口
# ---------------------------------------------------------------------------


def office_to_markdown(
    file_path: str,
    images_dir: str,
    rustfs_client=None,
) -> Tuple[str, List[Dict]]:
    """
    统一入口：根据文件扩展名自动选择对应转换器。

    Returns:
        (markdown_content, images)
        images: [{"src", "alt", "local_path", "rustfs_url"}]
    """
    ext = Path(file_path).suffix.lower()
    if ext == ".docx":
        return docx_to_markdown(file_path, images_dir, rustfs_client)
    elif ext == ".pptx":
        return pptx_to_markdown(file_path, images_dir, rustfs_client)
    elif ext == ".xlsx":
        return xlsx_to_markdown(file_path, images_dir, rustfs_client)
    else:
        raise ValueError(f"Unsupported format for office_to_markdown: {ext}")
